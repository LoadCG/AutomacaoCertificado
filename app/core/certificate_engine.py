"""
Motor de geração de certificados em lote — thread-safe.

Responsável por:
- Reabrir o template a cada certificado (nunca deepcopy — python-pptx usa ZIP)
- Normalizar runs quebrados antes de substituir
- Preservar formatação original de cada run (fonte, tamanho, negrito, itálico, cor)
- Publicar eventos de progresso na queue sem bloquear a UI
- Exportar .pptx e opcionalmente .pdf via PowerPoint COM (requer Office no Windows)
- Logar e pular erros individuais sem abortar o lote

Uso:
    from queue import Queue
    from app.core.certificate_engine import gerar_lote
    from app.utils.events import EventoGerador

    fila: Queue[EventoGerador] = Queue()
    gerar_lote(
        template=Path("template.pptx"),
        dados=df,
        mapeamento={"{{NOME}}": "Nome Completo", "{{RG}}": "RG"},
        pasta_saida=Path("saida/"),
        fila=fila,
        exportar_pdf=False,
    )
"""

import csv
import re
import shutil
import subprocess
import sys
import unicodedata
from datetime import date, datetime
from pathlib import Path
from queue import Queue
from typing import Optional

import pandas as pd
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.text.text import _Paragraph, _Run

from app.utils.events import (
    EventoConcluido,
    EventoErro,
    EventoGerador,
    EventoProgresso,
    EventoSucesso,
)
from app.utils.logger import obter_logger
from app.core.template_parser import _iterar_shapes

log = obter_logger(__name__)

# Regex para encontrar variáveis no formato {{VARIAVEL}}
PADRAO_VARIAVEL: re.Pattern = re.compile(r"\{\{[A-Z0-9_]+\}\}")

# Verifica disponibilidade do COM interop para exportação PDF
_COM_DISPONIVEL: bool = sys.platform == "win32"
try:
    import win32com.client  # type: ignore
    import pythoncom  # type: ignore
    _COM_DISPONIVEL = True
except ImportError:
    _COM_DISPONIVEL = False
    log.debug("pywin32 não disponível — exportação PDF desabilitada.")


# ---------------------------------------------------------------------------
# Exceções tipadas
# ---------------------------------------------------------------------------


class ExportacaoPDFError(Exception):
    """Levantada quando a exportação para PDF via COM interop falha."""


# ---------------------------------------------------------------------------
# Função principal — chamada em thread separada
# ---------------------------------------------------------------------------


# Padrão padrão de nome de arquivo: Nome do participante + data de hoje
PADRAO_NOME_PADRAO: str = "{{NOME}} - {DATA}"


def gerar_lote(
    template: Path,
    dados: pd.DataFrame,
    mapeamento: dict[str, str],
    pasta_saida: Path,
    fila: "Queue[EventoGerador]",
    exportar_pdf: bool = False,
    padrao_nome: str = PADRAO_NOME_PADRAO,
) -> None:
    """
    Gera certificados em lote a partir de um template .pptx e um DataFrame.

    Deve ser executada em uma thread separada — nunca na thread principal da UI.
    Publica eventos na fila para que a UI possa atualizar barra de progresso e log.

    Args:
        template: Caminho para o arquivo .pptx de template.
        dados: DataFrame com os dados dos participantes.
        mapeamento: Dicionário {variavel: nome_coluna}, ex: {"{{NOME}}": "Nome Completo"}.
        pasta_saida: Diretório onde os certificados serão salvos.
        fila: Queue de comunicação com a thread da UI.
        exportar_pdf: Se True, exporta também em .pdf via PowerPoint COM.
        padrao_nome: Padrão para o nome do arquivo, ex: '{{NOME}} - {DATA}'.\n
            Suporta variáveis do template ({{VAR}}) e variáveis de sistema:\n
            {DATA} = data atual (YYYY-MM-DD), {INDICE} = número da linha (0001...).

    Notes:
        - Reabre o template do disco a cada certificado (não usa deepcopy).
        - Erros individuais são logados e pulados — o lote continua.
        - Nomes de arquivo são sanitizados para remover caracteres inválidos em paths.
    """
    template = Path(template)
    ext_template = template.suffix.lower()
    if ext_template not in (".pptx", ".docx"):
        raise ValueError(f"Formato de template não suportado: {ext_template}")

    pasta_saida = Path(pasta_saida)
    pasta_saida.mkdir(parents=True, exist_ok=True)

    total = len(dados)
    total_sucesso = 0
    total_erro = 0

    log.info(
        "Iniciando geração de %d certificados. Template: %s | Saída: %s",
        total,
        template,
        pasta_saida,
    )

    for indice, (_, linha) in enumerate(dados.iterrows(), start=1):
        # Resolve os valores da linha para o mapeamento de variáveis
        mapeamento_valores = _resolver_valores(mapeamento, linha)

        # 1. Define o nome base pelo padrão do usuário
        nome_base = formatar_nome_arquivo(padrao_nome, mapeamento_valores, indice)
        
        # 2. Garante unicidade: evita que arquivos com mesmo nome se sobrescrevam
        caminho_gerado = pasta_saida / f"{nome_base}{ext_template}"
        
        contador_conflito = 1
        while caminho_gerado.exists():
            caminho_gerado = pasta_saida / f"{nome_base}-({contador_conflito}){ext_template}"
            contador_conflito += 1

        nome_final_sanitizado = caminho_gerado.stem

        try:
            if ext_template == ".pptx":
                # Reabre o template do disco — NUNCA deepcopy (python-pptx usa ZIP interno)
                prs = Presentation(str(template))
                for slide in prs.slides:
                    _processar_slide(slide, mapeamento_valores)
                prs.save(str(caminho_gerado))
                log.debug("Certificado PPTX salvo: %s", caminho_gerado)
            else:
                from docx import Document
                doc = Document(str(template))
                _processar_docx(doc, mapeamento_valores)
                doc.save(str(caminho_gerado))
                log.debug("Certificado DOCX salvo: %s", caminho_gerado)

            if exportar_pdf:
                caminho_pdf = pasta_saida / f"{nome_final_sanitizado}.pdf"
                sucesso_pdf = False

                # 1. Tenta exportar via COM (PowerPoint) se disponível (apenas para PPTX)
                if ext_template == ".pptx" and _COM_DISPONIVEL:
                    try:
                        _exportar_pdf_com(caminho_gerado, caminho_pdf)
                        log.debug("PDF exportado via PowerPoint COM: %s", caminho_pdf)
                        sucesso_pdf = True
                    except Exception as e_pdf:
                        log.warning(
                            "Falha ao exportar PDF via PowerPoint COM para '%s': %s. Tentando LibreOffice...",
                            nome_final_sanitizado,
                            e_pdf,
                        )

                # 2. Se falhar ou se for DOCX, tenta exportar via LibreOffice
                if not sucesso_pdf:
                    caminho_libreoffice = obter_caminho_libreoffice()
                    if caminho_libreoffice:
                        try:
                            _exportar_pdf_libreoffice(
                                caminho_libreoffice,
                                caminho_gerado,
                                pasta_saida,
                                caminho_pdf,
                            )
                            log.debug("PDF exportado via LibreOffice: %s", caminho_pdf)
                            sucesso_pdf = True
                        except Exception as e_lo:
                            log.warning(
                                "Falha ao exportar PDF via LibreOffice para '%s': %s",
                                nome_final_sanitizado,
                                e_lo,
                            )
                    else:
                        log.warning(
                            "Falha ao exportar PDF para '%s': PowerPoint COM falhou/indisponível/incompatível e LibreOffice não foi encontrado.",
                            nome_final_sanitizado,
                        )

            total_sucesso += 1
            fila.put(EventoSucesso(tipo="sucesso", arquivo=caminho_gerado.name))

        except Exception as e:
            total_erro += 1
            motivo = str(e)
            log.error(
                "Erro ao gerar certificado '%s' (linha %d): %s",
                nome_final_sanitizado,
                indice,
                motivo,
                exc_info=True,
            )
            fila.put(
                EventoErro(
                    tipo="erro",
                    linha=indice,
                    arquivo=f"{nome_final_sanitizado}{ext_template}",
                    motivo=motivo,
                )
            )

        # Publica progresso após cada item (sucesso ou erro)
        fila.put(EventoProgresso(tipo="progresso", atual=indice, total=total))

    fila.put(
        EventoConcluido(
            tipo="concluido",
            total_sucesso=total_sucesso,
            total_erro=total_erro,
        )
    )
    log.info(
        "Lote concluído: %d sucesso, %d erro(s).", total_sucesso, total_erro
    )


# ---------------------------------------------------------------------------
# Processamento de slides e substituição
# ---------------------------------------------------------------------------


def _processar_slide(slide, mapeamento_valores: dict[str, str]) -> None:
    """
    Processa todos os parágrafos de um slide, normalizando runs e
    substituindo variáveis.

    Args:
        slide: Objeto Slide do python-pptx.
        mapeamento_valores: Dicionário {variavel: valor}, ex: {"{{NOME}}": "João"}.
    """
    for shape in _iterar_shapes(slide):
        if not shape.has_text_frame:
            continue
        for paragrafo in shape.text_frame.paragraphs:
            # Normaliza runs quebrados ANTES de tentar substituir
            _normalizar_runs_paragrafo(paragrafo)
            # Substitui variáveis em cada run preservando formatação
            for run in paragrafo.runs:
                for variavel, valor in mapeamento_valores.items():
                    _substituir_run(run, variavel, valor)


def _processar_docx(doc, mapeamento_valores: dict[str, str]) -> None:
    """
    Processa todos os parágrafos e tabelas de um documento Word (.docx),
    substituindo as variáveis configuradas.
    """
    # 1. Processa parágrafos normais
    for p in doc.paragraphs:
        _substituir_no_paragrafo_docx(p, mapeamento_valores)

    # 2. Processa tabelas
    for tabela in doc.tables:
        for linha in tabela.rows:
            for celula in linha.cells:
                for p in celula.paragraphs:
                    _substituir_no_paragrafo_docx(p, mapeamento_valores)


def _substituir_no_paragrafo_docx(p, mapeamento_valores: dict[str, str]) -> None:
    """
    Substitui as variáveis no texto de um parágrafo do Word preservando
    a formatação básica do primeiro run se houver alteração.
    """
    texto_antigo = p.text
    texto_novo = texto_antigo
    for var, valor in mapeamento_valores.items():
        texto_novo = texto_novo.replace(var, valor)

    if texto_novo == texto_antigo:
        return

    if p.runs:
        primeiro_run = p.runs[0]
        bold = primeiro_run.bold
        italic = primeiro_run.italic
        
        font_name = None
        font_size = None
        font_color = None
        if hasattr(primeiro_run, "font") and primeiro_run.font:
            font_name = primeiro_run.font.name
            font_size = primeiro_run.font.size
            if primeiro_run.font.color:
                try:
                    font_color = primeiro_run.font.color.rgb
                except Exception:
                    pass

        p.text = ""
        novo_run = p.add_run(texto_novo)
        novo_run.bold = bold
        novo_run.italic = italic
        if font_name:
            novo_run.font.name = font_name
        if font_size:
            novo_run.font.size = font_size
        if font_color:
            try:
                novo_run.font.color.rgb = font_color
            except Exception:
                pass
    else:
        p.text = texto_novo


def _normalizar_runs_paragrafo(paragrafo: _Paragraph) -> None:
    """
    Mescla runs consecutivos quando uma variável está quebrada entre eles.

    O PowerPoint frequentemente divide '{{NOME}}' em múltiplos runs ao
    editar texto formatado: ['{{NO', 'ME}}']. Esta função detecta essa
    situação e mescla todos os runs no primeiro, zerando o texto dos demais.

    O primeiro run preserva sua formatação original. O comportamento
    é: se o parágrafo contém uma variável quando o texto é concatenado,
    mas nenhum run individual a contém completamente — mesclar tudo.

    Args:
        paragrafo: Objeto _Paragraph do python-pptx a ser normalizado.
    """
    runs = paragrafo.runs
    if len(runs) <= 1:
        return

    texto_completo = "".join(r.text for r in runs)

    # Verifica se há variáveis no texto completo mas quebradas nos runs
    variaveis_no_total = set(PADRAO_VARIAVEL.findall(texto_completo))
    if not variaveis_no_total:
        return  # Nenhuma variável — não precisa normalizar

    variaveis_em_runs = set()
    for run in runs:
        variaveis_em_runs.update(PADRAO_VARIAVEL.findall(run.text))

    # Se todas as variáveis já estão inteiras em algum run, não precisa mesclar
    if variaveis_no_total == variaveis_em_runs:
        return

    # Mescla tudo no primeiro run e zera os demais
    runs[0].text = texto_completo
    for run in runs[1:]:
        run.text = ""

    log.debug(
        "Runs normalizados (texto mesclado no run[0]): '%s'", texto_completo[:80]
    )


def _substituir_run(run: _Run, variavel: str, valor: str) -> None:
    """
    Substitui uma variável no texto do run preservando toda a formatação.

    Faz um snapshot da formatação antes de alterar o texto e restaura
    após — necessário porque atribuir `run.text` reseta a formatação
    para o estilo padrão do parágrafo.

    Verifica `MSO_THEME_COLOR` antes de acessar `.rgb` para evitar
    AttributeError em runs sem cor explícita.

    Args:
        run: Objeto _Run do python-pptx.
        variavel: Variável a substituir, ex: '{{NOME}}'.
        valor: Valor a inserir no lugar da variável.
    """
    if variavel not in run.text:
        return

    # Snapshot da formatação antes de qualquer modificação
    fonte = run.font
    snap_bold = fonte.bold
    snap_italic = fonte.italic
    snap_size = fonte.size
    snap_name = fonte.name
    snap_underline = fonte.underline

    # Cor: verificar type antes de acessar .rgb (runs sem cor = type None)
    snap_cor: Optional[RGBColor] = None
    try:
        from pptx.enum.dml import MSO_COLOR_TYPE
        if fonte.color.type == MSO_COLOR_TYPE.RGB:
            snap_cor = fonte.color.rgb
    except (AttributeError, TypeError):
        pass  # Run sem cor explícita — ignorar

    # Realiza a substituição (isso pode resetar formatação internamente)
    run.text = run.text.replace(variavel, str(valor))

    # Restaura a formatação do snapshot
    fonte.bold = snap_bold
    fonte.italic = snap_italic
    fonte.size = snap_size
    fonte.underline = snap_underline

    if snap_name:
        fonte.name = snap_name

    if snap_cor is not None:
        try:
            fonte.color.rgb = snap_cor
        except AttributeError:
            pass  # Ignorar falha ao restaurar cor


# ---------------------------------------------------------------------------
# Exportação PDF via COM interop
# ---------------------------------------------------------------------------


def _exportar_pdf_com(caminho_pptx: Path, caminho_pdf: Path) -> None:
    """
    Exporta um arquivo .pptx para .pdf via PowerPoint COM interop.

    Requer Microsoft Office instalado no sistema. Verifica se o PowerPoint
    já estava aberto para não encerrá-lo ao terminar (comportamento seguro).

    Args:
        caminho_pptx: Caminho absoluto do arquivo .pptx de entrada.
        caminho_pdf: Caminho absoluto do arquivo .pdf de saída.

    Raises:
        ExportacaoPDFError: Se o PowerPoint não estiver disponível ou a
            exportação falhar por qualquer motivo.
    """
    if not _COM_DISPONIVEL:
        raise ExportacaoPDFError(
            "pywin32 não está instalado. Execute: pip install pywin32"
        )

    pptx_abs = str(caminho_pptx.resolve())
    pdf_abs = str(caminho_pdf.resolve())

    ppt = None
    ppt_ja_estava_aberto = False
    com_inicializado = False

    try:
        # CoInitialize deve ser chamado uma vez por thread — ignora se já inicializado
        try:
            pythoncom.CoInitialize()  # type: ignore
            com_inicializado = True
        except Exception:
            pass  # Já inicializado nesta thread — tudo bem

        # Tenta se conectar a uma instância já aberta antes de criar nova
        # Isso evita fechar o PowerPoint do usuário ao terminar
        try:
            ppt = win32com.client.GetActiveObject("PowerPoint.Application")  # type: ignore
            ppt_ja_estava_aberto = True
            log.debug("Usando instância do PowerPoint já aberta.")
        except Exception:
            ppt = win32com.client.Dispatch("PowerPoint.Application")  # type: ignore
            ppt_ja_estava_aberto = False

        apresentacao = ppt.Presentations.Open(pptx_abs, WithWindow=False)
        try:
            apresentacao.SaveAs(pdf_abs, 32)  # 32 = ppSaveAsPDF
        finally:
            apresentacao.Close()

    except Exception as e:
        raise ExportacaoPDFError(
            f"Falha ao exportar PDF '{caminho_pdf.name}': {e}\n"
            "Verifique se o Microsoft Office está instalado e tente novamente."
        ) from e
    finally:
        # Só encerra o PowerPoint se fomos nós que o abrimos
        if ppt is not None and not ppt_ja_estava_aberto:
            try:
                ppt.Quit()
            except Exception:
                pass
        if com_inicializado:
            try:
                pythoncom.CoUninitialize()  # type: ignore
            except Exception:
                pass


def com_disponivel() -> bool:
    """
    Informa se a exportação PDF via COM interop está disponível.

    Returns:
        True se pywin32 estiver instalado e o sistema for Windows.
    """
    return _COM_DISPONIVEL


def obter_caminho_libreoffice() -> Optional[Path]:
    """
    Tenta localizar o executável do LibreOffice (soffice) no sistema.

    Procura no PATH e em caminhos padrão de instalação do Windows.
    """
    # 1. Tenta encontrar no PATH
    soffice_path = shutil.which("soffice")
    if soffice_path:
        return Path(soffice_path)

    # 2. Caminhos padrão no Windows
    caminhos_padrao = [
        Path(r"C:\Program Files\LibreOffice\program\soffice.exe"),
        Path(r"C:\Program Files (x86)\LibreOffice\program\soffice.exe"),
    ]
    for caminho in caminhos_padrao:
        if caminho.is_file():
            return caminho

    return None


def pdf_disponivel() -> bool:
    """
    Informa se a exportação PDF está disponível no sistema por algum meio.

    Returns:
        True se PowerPoint COM estiver disponível ou se LibreOffice estiver instalado.
    """
    return _COM_DISPONIVEL or (obter_caminho_libreoffice() is not None)


def _exportar_pdf_libreoffice(
    caminho_libreoffice: Path,
    caminho_pptx: Path,
    pasta_saida: Path,
    caminho_pdf: Path,
) -> None:
    """
    Exporta um arquivo .pptx para .pdf usando o LibreOffice em modo headless.

    Args:
        caminho_libreoffice: Caminho absoluto para o executável soffice.
        caminho_pptx: Caminho do arquivo de apresentação de entrada.
        pasta_saida: Diretório onde o PDF gerado será salvo inicialmente.
        caminho_pdf: Caminho desejado para o PDF final (caso precise de renomeação).
    """
    cmd = [
        str(caminho_libreoffice),
        "--headless",
        "--convert-to",
        "pdf",
        "--outdir",
        str(pasta_saida.resolve()),
        str(caminho_pptx.resolve()),
    ]
    result = subprocess.run(
        cmd, stdout=subprocess.PIPE, stderr=subprocess.PIPE, text=True
    )
    if result.returncode != 0:
        raise RuntimeError(
            f"Erro na conversão via LibreOffice (Código {result.returncode}): {result.stderr}"
        )

    # O LibreOffice gera o PDF na pasta de saída com o mesmo nome do slide.
    # Exemplo: pasta_saida / "Nome-do-Slide.pdf"
    pdf_gerado = pasta_saida / f"{caminho_pptx.stem}.pdf"
    if pdf_gerado.exists():
        if pdf_gerado.resolve() != caminho_pdf.resolve():
            if caminho_pdf.exists():
                caminho_pdf.unlink()
            pdf_gerado.rename(caminho_pdf)
    else:
        raise RuntimeError(
            f"O arquivo PDF esperado não foi gerado em: {pdf_gerado}"
        )


def escrever_relatorio_erros(
    pasta_saida: Path, erros: list["EventoErro"]
) -> Optional[Path]:
    """
    Escreve um relatório CSV com os erros ocorridos durante um lote.

    Args:
        pasta_saida: Diretório onde o relatório será salvo.
        erros: Lista de eventos de erro coletados durante a geração.

    Returns:
        Caminho do arquivo de relatório gerado, ou None se `erros` estiver vazio.
    """
    if not erros:
        return None

    pasta_saida = Path(pasta_saida)
    timestamp = datetime.now().strftime("%Y%m%d-%H%M%S")
    caminho_relatorio = pasta_saida / f"relatorio-erros-{timestamp}.csv"

    with open(caminho_relatorio, "w", newline="", encoding="utf-8-sig") as arquivo:
        escritor = csv.writer(arquivo, delimiter=";")
        escritor.writerow(["Linha", "Arquivo", "Motivo"])
        for erro in erros:
            escritor.writerow([erro["linha"], erro["arquivo"], erro["motivo"]])

    log.info("Relatório de erros salvo: %s (%d erro(s))", caminho_relatorio, len(erros))
    return caminho_relatorio


# ---------------------------------------------------------------------------
# Funções auxiliares
# ---------------------------------------------------------------------------


def _resolver_valores(
    mapeamento: dict[str, str], linha: pd.Series
) -> dict[str, str]:
    """
    Resolve o mapeamento variável→coluna para variável→valor usando a linha atual.

    Suporta:
    1. Mapeamento para coluna: valor buscado na linha da planilha.
    2. Mapeamento para texto fixo: prefixo 'FIXED:' indica valor literal.

    Args:
        mapeamento: Dicionário {variavel: nome_coluna_ou_texto}.
        linha: Série pandas com os dados de um participante.

    Returns:
        Dicionário {variavel: valor_str} pronto para substituição.
    """
    resultado: dict[str, str] = {}
    for variavel, mapeado in mapeamento.items():
        # Caso 1: Texto fixo (prefixo FIXED:)
        if mapeado.startswith("FIXED:"):
            resultado[variavel] = mapeado[6:].strip()
            continue

        # Caso 2: Coluna da planilha
        nome_coluna = mapeado
        if nome_coluna not in linha.index:
            log.warning(
                "Coluna '%s' não encontrada na linha. Usando string vazia.", nome_coluna
            )
            resultado[variavel] = ""
            continue

        valor = linha[nome_coluna]
        # Converte NaN e None para string vazia
        if pd.isna(valor) if not isinstance(valor, str) else False:
            resultado[variavel] = ""
        else:
            resultado[variavel] = str(valor).strip()
    return resultado


def formatar_nome_arquivo(
    padrao: str,
    mapeamento_valores: dict[str, str],
    indice: int,
) -> str:
    """
    Gera o nome do arquivo de saída usando o padrão configurado pelo usuário.

    Substitui:
    - {{VARIAVEL}} → valor da variável do template para aquele participante
    - {DATA}       → data atual no formato YYYY-MM-DD
    - {INDICE}     → número da linha com 4 dígitos (0001, 0002...)

    Variáveis não substituídas são removidas silenciosamente.
    O resultado é sanitizado para uso como nome de arquivo no Windows.

    Args:
        padrao: Padrão de nome, ex: '{{NOME}} - {DATA}' ou 'Cert_{{NOME}}_{{TURMA}}'.
        mapeamento_valores: Dicionário {variavel: valor} já resolvido para a linha atual.
        indice: Posição da linha no DataFrame (base 1).

    Returns:
        String sanitizada com até 200 caracteres.
    """
    nome = padrao

    # Substitui variáveis do template ({{VAR}})
    for variavel, valor in mapeamento_valores.items():
        nome = nome.replace(variavel, str(valor))

    # Substitui variáveis de sistema
    nome = nome.replace("{DATA}", date.today().strftime("%Y-%m-%d"))
    nome = nome.replace("{INDICE}", f"{indice:04d}")

    # Remove variáveis não substituídas ({{VAR}} ou {VAR} que sobraram)
    nome = re.sub(r"\{\{[A-Z0-9_]+\}\}", "", nome)
    nome = re.sub(r"\{[A-Z_]+\}", "", nome)

    # Remove caracteres inválidos em paths do Windows: \ / : * ? " < > |
    nome = re.sub(r'[\\/:*?"<>|]', "_", nome)

    # Padronização: Espaços e sublinhados viram hifens
    nome = nome.replace(" ", "-").replace("_", "-")

    # Limpa hifens múltiplos (ex: "---" -> "-") e espaços que sobraram
    nome = re.sub(r"[-]{2,}", "-", nome)
    nome = nome.strip("-")

    # Fallback se o padrão resultou em string vazia
    if not nome.strip():
        nome = f"Certificado-{indice:04d}"

    return nome[:200]
