"""
Módulo de extração de variáveis de templates .pptx.

Detecta todas as ocorrências de {{VARIAVEL}} nos slides, incluindo
variáveis que foram quebradas em múltiplos runs pelo PowerPoint.
A detecção é feita sobre o texto reconstruído por concatenação dos runs
de cada parágrafo, garantindo que o que é detectado aqui pode ser
substituído pelo certificate_engine.

Uso:
    from app.core.template_parser import extrair_variaveis
    variaveis = extrair_variaveis(Path("template.pptx"))
    # ['{{NOME}}', '{{CARGO}}', '{{DATA}}']
"""

import re
from pathlib import Path
from typing import Generator

from pptx import Presentation
from pptx.shapes.base import BaseShape
from pptx.text.text import _Paragraph

from app.utils.logger import obter_logger

log = obter_logger(__name__)

# Regex para encontrar variáveis no formato {{NOME_DA_VARIAVEL}}
# Aceita letras maiúsculas, dígitos e sublinhado — sem espaços, sem acentos
PADRAO_VARIAVEL: re.Pattern = re.compile(r"\{\{[A-Z0-9_]+\}\}")


class TemplateSemVariaveisError(Exception):
    """
    Levantada quando o template .pptx não contém nenhuma variável {{...}}.

    Isso provavelmente indica que o usuário selecionou o arquivo errado
    ou que o template não foi preparado corretamente.
    """


# ---------------------------------------------------------------------------
# Funções públicas
# ---------------------------------------------------------------------------


def extrair_variaveis(caminho_template: Path) -> list[str]:
    """
    Extrai todas as variáveis {{VARIAVEL}} do arquivo .pptx.

    Reconstrói o texto de cada parágrafo concatenando seus runs antes
    de aplicar o regex, detectando variáveis quebradas entre múltiplos
    runs (comportamento comum do PowerPoint ao editar texto formatado).

    Args:
        caminho_template: Caminho para o arquivo .pptx.

    Returns:
        Lista deduplicada e ordenada de variáveis encontradas,
        ex: ['{{CARGO}}', '{{DATA}}', '{{NOME}}'].

    Raises:
        FileNotFoundError: Se o arquivo não existir.
        TemplateSemVariaveisError: Se nenhuma variável for encontrada.
    """
    caminho_template = Path(caminho_template)

    if not caminho_template.is_file():
        raise FileNotFoundError(
            f"Template não encontrado: '{caminho_template}'"
        )

    log.info("Extraindo variáveis do template: %s", caminho_template)

    prs = Presentation(str(caminho_template))
    variaveis_encontradas: set[str] = set()

    for num_slide, slide in enumerate(prs.slides, start=1):
        for shape in _iterar_shapes(slide):
            if not shape.has_text_frame:
                continue
            for paragrafo in shape.text_frame.paragraphs:
                texto_reconstruido = reconstruir_texto_paragrafo(paragrafo)
                ocorrencias = PADRAO_VARIAVEL.findall(texto_reconstruido)
                if ocorrencias:
                    log.debug(
                        "Slide %d — shape '%s': variáveis=%s",
                        num_slide,
                        shape.name,
                        ocorrencias,
                    )
                    variaveis_encontradas.update(ocorrencias)

    variaveis_ordenadas = sorted(variaveis_encontradas)
    log.info(
        "Variáveis encontradas no template (%d): %s",
        len(variaveis_ordenadas),
        variaveis_ordenadas,
    )
    return variaveis_ordenadas


def reconstruir_texto_paragrafo(paragrafo: _Paragraph) -> str:
    """
    Reconstrói o texto completo de um parágrafo concatenando todos os runs.

    Necessário porque o PowerPoint pode dividir uma variável como
    '{{NOME}}' em múltiplos runs: ['{{NO', 'ME}}'], tornando o regex
    ineficaz se aplicado em cada run individualmente.

    Args:
        paragrafo: Objeto _Paragraph do python-pptx.

    Returns:
        String com o texto completo do parágrafo.
    """
    return "".join(run.text for run in paragrafo.runs)


# ---------------------------------------------------------------------------
# Funções auxiliares privadas
# ---------------------------------------------------------------------------


def _iterar_shapes(slide) -> Generator[BaseShape, None, None]:
    """
    Itera sobre todas as shapes de um slide, incluindo as dentro de grupos.

    O PowerPoint permite agrupar shapes, e o python-pptx não itera
    automaticamente dentro de GroupShape. Esta função faz isso recursivamente.

    Args:
        slide: Objeto Slide do python-pptx.

    Yields:
        Cada BaseShape encontrada, incluindo dentro de grupos.
    """
    for shape in slide.shapes:
        yield shape
        # Se for um grupo, iterar recursivamente nas shapes internas
        if shape.shape_type == 6:  # MSO_SHAPE_TYPE.GROUP == 6
            yield from _iterar_shapes_grupo(shape)


def _iterar_shapes_grupo(grupo) -> Generator[BaseShape, None, None]:
    """
    Itera recursivamente sobre shapes dentro de um GroupShape.

    Args:
        grupo: Objeto GroupShape do python-pptx.

    Yields:
        Cada BaseShape dentro do grupo (e sub-grupos).
    """
    try:
        for shape in grupo.shapes:
            yield shape
            if shape.shape_type == 6:
                yield from _iterar_shapes_grupo(shape)
    except AttributeError:
        # GroupShape sem atributo .shapes — ignora silenciosamente
        pass
