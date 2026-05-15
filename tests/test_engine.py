"""Testes unitários e de integração do certificate_engine."""

import time
from pathlib import Path
from queue import Queue

import pytest
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

from app.core.certificate_engine import (
    _normalizar_runs_paragrafo,
    _resolver_valores,
    _gerar_nome_por_padrao,
    _substituir_run,
    gerar_lote,
    PADRAO_NOME_PADRAO,
)
from app.utils.events import EventoGerador


# ---------------------------------------------------------------------------
# Testes de _substituir_run
# ---------------------------------------------------------------------------


class TestSubstituirRun:
    def _criar_run(self, texto: str, bold=None, italic=None, size=None, cor=None):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
        run = box.text_frame.paragraphs[0].add_run()
        run.text = texto
        if bold is not None:
            run.font.bold = bold
        if italic is not None:
            run.font.italic = italic
        if size is not None:
            run.font.size = size
        if cor is not None:
            run.font.color.rgb = cor
        return run

    def test_substituicao_basica(self):
        run = self._criar_run("{{NOME}}")
        _substituir_run(run, "{{NOME}}", "Alice")
        assert run.text == "Alice"

    def test_preserva_bold(self):
        run = self._criar_run("{{NOME}}", bold=True)
        _substituir_run(run, "{{NOME}}", "Alice")
        assert run.font.bold is True

    def test_preserva_italic(self):
        run = self._criar_run("{{NOME}}", italic=True)
        _substituir_run(run, "{{NOME}}", "Alice")
        assert run.font.italic is True

    def test_preserva_tamanho(self):
        run = self._criar_run("{{NOME}}", size=Pt(24))
        _substituir_run(run, "{{NOME}}", "Alice")
        assert run.font.size == Pt(24)

    def test_preserva_cor_rgb(self):
        cor = RGBColor(0x1E, 0x90, 0xFF)
        run = self._criar_run("{{NOME}}", cor=cor)
        _substituir_run(run, "{{NOME}}", "Alice")
        assert run.font.color.rgb == cor

    def test_run_sem_cor_explicita_nao_lanca_erro(self):
        """
        Run sem cor explícita (font.color.type = None) não deve lançar
        AttributeError — cobertura do bug corrigido no plano revisado.
        """
        run = self._criar_run("{{NOME}}")
        # Sem definir cor — simula run sem formatação explícita de cor
        _substituir_run(run, "{{NOME}}", "Maria Silva")
        assert run.text == "Maria Silva"

    def test_chave_ausente_nao_modifica(self):
        run = self._criar_run("Texto fixo")
        _substituir_run(run, "{{NOME}}", "Alice")
        assert run.text == "Texto fixo"

    def test_valor_com_caracteres_especiais(self):
        run = self._criar_run("{{NOME}}")
        _substituir_run(run, "{{NOME}}", "José Ação Ângela")
        assert run.text == "José Ação Ângela"


class TestResolverValores:
    def test_resolve_coluna(self, df_simples):
        linha = df_simples.iloc[0]
        mapeamento = {"{{NOME}}": "Nome", "{{RG}}": "RG"}
        resultado = _resolver_valores(mapeamento, linha)
        assert resultado["{{NOME}}"] == "Alice Silva"
        assert resultado["{{RG}}"] == "12.345.678-9"

    def test_resolve_texto_fixo(self, df_simples):
        linha = df_simples.iloc[0]
        mapeamento = {"{{NOME}}": "FIXED:Texto Literal", "{{DATA}}": "FIXED: 2024 "}
        resultado = _resolver_valores(mapeamento, linha)
        assert resultado["{{NOME}}"] == "Texto Literal"
        assert resultado["{{DATA}}"] == "2024"

    def test_resolve_misto(self, df_simples):
        linha = df_simples.iloc[0]
        mapeamento = {"{{NOME}}": "Nome", "{{EVENTO}}": "FIXED:Workshop"}
        resultado = _resolver_valores(mapeamento, linha)
        assert resultado["{{NOME}}"] == "Alice Silva"
        assert resultado["{{EVENTO}}"] == "Workshop"


# ---------------------------------------------------------------------------
# Testes de _normalizar_runs_paragrafo
# ---------------------------------------------------------------------------


class TestNormalizarRuns:
    def test_normaliza_run_quebrado(self, template_run_quebrado: Path):
        prs = Presentation(str(template_run_quebrado))
        paragrafo = prs.slides[0].shapes[0].text_frame.paragraphs[0]
        assert paragrafo.runs[0].text == "{{NO"
        _normalizar_runs_paragrafo(paragrafo)
        assert "{{NOME}}" in paragrafo.runs[0].text

    def test_nao_modifica_run_integro(self, template_simples: Path):
        prs = Presentation(str(template_simples))
        paragrafo = prs.slides[0].shapes[0].text_frame.paragraphs[0]
        texto_antes = paragrafo.runs[0].text
        _normalizar_runs_paragrafo(paragrafo)
        assert paragrafo.runs[0].text == texto_antes


# ---------------------------------------------------------------------------
# Testes de _gerar_nome_arquivo
# ---------------------------------------------------------------------------


class TestGerarNomePorPadrao:
    def test_padrao_padrao_contem_nome(self):
        resultado = _gerar_nome_por_padrao(
            PADRAO_NOME_PADRAO, {"{{NOME}}": "Alice Silva"}, 1
        )
        assert "Alice Silva" in resultado

    def test_padrao_padrao_contem_data(self):
        from datetime import date
        resultado = _gerar_nome_por_padrao(
            PADRAO_NOME_PADRAO, {"{{NOME}}": "Alice"}, 1
        )
        assert date.today().strftime("%Y-%m-%d") in resultado

    def test_padrao_customizado(self):
        resultado = _gerar_nome_por_padrao(
            "Cert_{{NOME}}_{INDICE}", {"{{NOME}}": "Bob"}, 7
        )
        assert "Cert_Bob" in resultado
        assert "0007" in resultado

    def test_sanitiza_barra_no_valor(self):
        resultado = _gerar_nome_por_padrao(
            "{{NOME}}", {"{{NOME}}": "Alice/Santos"}, 1
        )
        assert "/" not in resultado

    def test_fallback_padrao_vazio(self):
        resultado = _gerar_nome_por_padrao("", {}, 42)
        assert "42" in resultado

    def test_variaveis_nao_mapeadas_removidas(self):
        resultado = _gerar_nome_por_padrao(
            "{{NOME}} - {{CARGO}} - {DATA}",
            {"{{NOME}}": "Alice"},  # {{CARGO}} não mapeado
            1,
        )
        assert "{{CARGO}}" not in resultado
        assert "Alice" in resultado


# ---------------------------------------------------------------------------
# Testes de integração — gerar_lote
# ---------------------------------------------------------------------------


class TestGerarLote:
    def test_gera_arquivos_corretos(
        self, template_simples: Path, df_simples, tmp_path: Path
    ):
        fila: Queue[EventoGerador] = Queue()
        mapeamento = {"{{NOME}}": "Nome", "{{RG}}": "RG"}
        pasta_saida = tmp_path / "saida"
        gerar_lote(template_simples, df_simples, mapeamento, pasta_saida, fila)

        pptx_gerados = list(pasta_saida.glob("*.pptx"))
        assert len(pptx_gerados) == 3

    def test_conteudo_substituido(
        self, template_simples: Path, df_simples, tmp_path: Path
    ):
        fila: Queue[EventoGerador] = Queue()
        mapeamento = {"{{NOME}}": "Nome", "{{RG}}": "RG"}
        pasta_saida = tmp_path / "saida2"
        gerar_lote(template_simples, df_simples, mapeamento, pasta_saida, fila)

        arquivo = next(pasta_saida.glob("*Alice*.pptx"), None)
        assert arquivo is not None
        prs = Presentation(str(arquivo))
        textos = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        textos.append("".join(r.text for r in para.runs))
        texto_completo = " ".join(textos)
        assert "Alice Silva" in texto_completo
        assert "{{NOME}}" not in texto_completo

    def test_evento_concluido_publicado(
        self, template_simples: Path, df_simples, tmp_path: Path
    ):
        fila: Queue[EventoGerador] = Queue()
        mapeamento = {"{{NOME}}": "Nome", "{{RG}}": "RG"}
        gerar_lote(template_simples, df_simples, mapeamento, tmp_path, fila)

        eventos = []
        while not fila.empty():
            eventos.append(fila.get_nowait())

        tipos = [e["tipo"] for e in eventos]
        assert "concluido" in tipos

        conclusao = next(e for e in eventos if e["tipo"] == "concluido")
        assert conclusao["total_sucesso"] == 3
        assert conclusao["total_erro"] == 0

    def test_erro_individual_nao_aborta_lote(
        self, template_simples: Path, df_com_nulos, tmp_path: Path
    ):
        """
        Linha com nome contendo '/' (inválido em path no Windows) causa
        erro no arquivo mas não aborta o lote.
        """
        fila: Queue[EventoGerador] = Queue()
        mapeamento = {"{{NOME}}": "Nome", "{{RG}}": "RG"}
        # df_com_nulos tem 3 linhas — mesmo com erro em alguma, deve continuar
        gerar_lote(template_simples, df_com_nulos, mapeamento, tmp_path, fila)

        eventos = []
        while not fila.empty():
            eventos.append(fila.get_nowait())

        conclusao = next((e for e in eventos if e["tipo"] == "concluido"), None)
        assert conclusao is not None
        # O lote deve ter terminado (sucesso + erro = total de linhas)
        total = conclusao["total_sucesso"] + conclusao["total_erro"]
        assert total == len(df_com_nulos)

    def test_nao_modifica_template_original(
        self, template_simples: Path, df_simples, tmp_path: Path
    ):
        """Garante que o template original não é modificado após a geração."""
        import os
        mtime_antes = os.path.getmtime(template_simples)
        fila: Queue[EventoGerador] = Queue()
        gerar_lote(
            template_simples, df_simples,
            {"{{NOME}}": "Nome", "{{RG}}": "RG"},
            tmp_path, fila
        )
        mtime_depois = os.path.getmtime(template_simples)
        assert mtime_antes == mtime_depois
