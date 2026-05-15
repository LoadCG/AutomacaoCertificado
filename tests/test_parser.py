"""Testes unitários do módulo template_parser."""

import pytest
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches

from app.core.template_parser import (
    extrair_variaveis,
    reconstruir_texto_paragrafo,
)


class TestExtrairVariaveis:
    def test_variaveis_simples(self, template_simples: Path):
        variaveis = extrair_variaveis(template_simples)
        assert "{{NOME}}" in variaveis
        assert "{{RG}}" in variaveis

    def test_variavel_quebrada_entre_dois_runs(self, template_run_quebrado: Path):
        """Variável dividida em runs deve ser detectada via texto reconstruído."""
        variaveis = extrair_variaveis(template_run_quebrado)
        assert "{{NOME}}" in variaveis

    def test_sem_variaveis_retorna_lista_vazia(self, template_sem_variaveis: Path):
        variaveis = extrair_variaveis(template_sem_variaveis)
        assert variaveis == []

    def test_variaveis_deduplicadas(self, tmp_path: Path):
        """Variável repetida em múltiplos slides deve aparecer uma única vez."""
        prs = Presentation()
        for _ in range(3):
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
            box.text_frame.paragraphs[0].add_run().text = "{{NOME}}"
        caminho = tmp_path / "dup.pptx"
        prs.save(str(caminho))
        variaveis = extrair_variaveis(caminho)
        assert variaveis.count("{{NOME}}") == 1

    def test_retorna_lista_ordenada(self, tmp_path: Path):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(2))
        box.text_frame.paragraphs[0].add_run().text = "{{ZULU}} {{ALPHA}} {{MIKE}}"
        caminho = tmp_path / "ordem.pptx"
        prs.save(str(caminho))
        variaveis = extrair_variaveis(caminho)
        assert variaveis == sorted(variaveis)

    def test_arquivo_inexistente(self, tmp_path: Path):
        with pytest.raises(FileNotFoundError):
            extrair_variaveis(tmp_path / "nao_existe.pptx")

    def test_variavel_case_sensitive(self, tmp_path: Path):
        """Variáveis em minúsculo NÃO devem ser detectadas (só uppercase)."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
        box.text_frame.paragraphs[0].add_run().text = "{{nome}}"
        caminho = tmp_path / "minusculo.pptx"
        prs.save(str(caminho))
        variaveis = extrair_variaveis(caminho)
        assert "{{nome}}" not in variaveis


class TestReconstruirTextoParagrafo:
    def test_paragrafo_run_unico(self, template_simples: Path):
        prs = Presentation(str(template_simples))
        paragrafo = prs.slides[0].shapes[0].text_frame.paragraphs[0]
        texto = reconstruir_texto_paragrafo(paragrafo)
        assert "{{NOME}}" in texto

    def test_paragrafo_multiplos_runs(self, template_run_quebrado: Path):
        prs = Presentation(str(template_run_quebrado))
        paragrafo = prs.slides[0].shapes[0].text_frame.paragraphs[0]
        texto = reconstruir_texto_paragrafo(paragrafo)
        assert texto == "{{NOME}}"
