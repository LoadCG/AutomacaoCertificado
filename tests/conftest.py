"""Fixtures compartilhadas para os testes do Gerador de Certificados."""

import io
import pytest
import pandas as pd
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor


@pytest.fixture
def df_simples() -> pd.DataFrame:
    """DataFrame com 3 participantes para testes de geração."""
    return pd.DataFrame({
        "Nome": ["Alice Silva", "Bob Santos", "Carol Mendes"],
        "RG": ["12.345.678-9", "98.765.432-1", "11.111.111-1"],
        "Cargo": ["Analista", "Gerente", "Coordenadora"],
    })


@pytest.fixture
def df_com_nulos() -> pd.DataFrame:
    """DataFrame com células vazias e caracteres especiais."""
    return pd.DataFrame({
        "Nome": ["José Ação", None, "Ângela/Souza"],
        "RG": ["111", "222", "333"],
    })


@pytest.fixture
def template_simples(tmp_path: Path) -> Path:
    """Template .pptx com variáveis {{NOME}} e {{RG}}."""
    prs = Presentation()
    layout = prs.slide_layouts[6]  # Slide em branco
    slide = prs.slides.add_slide(layout)

    txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1))
    tf = txBox.text_frame
    run = tf.paragraphs[0].add_run()
    run.text = "{{NOME}}"
    run.font.bold = True
    run.font.size = Pt(24)
    run.font.color.rgb = RGBColor(0x1E, 0x90, 0xFF)

    txBox2 = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(5), Inches(1))
    run2 = txBox2.text_frame.paragraphs[0].add_run()
    run2.text = "{{RG}}"

    caminho = tmp_path / "template_test.pptx"
    prs.save(str(caminho))
    return caminho


@pytest.fixture
def template_run_quebrado(tmp_path: Path) -> Path:
    """Template com variável {{NOME}} dividida entre 2 runs."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1))
    tf = txBox.text_frame
    paragrafo = tf.paragraphs[0]

    run1 = paragrafo.add_run()
    run1.text = "{{NO"

    run2 = paragrafo.add_run()
    run2.text = "ME}}"

    caminho = tmp_path / "template_quebrado.pptx"
    prs.save(str(caminho))
    return caminho


@pytest.fixture
def template_sem_variaveis(tmp_path: Path) -> Path:
    """Template sem nenhuma variável."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1))
    txBox.text_frame.paragraphs[0].add_run().text = "Texto fixo sem variável"
    caminho = tmp_path / "template_sem_vars.pptx"
    prs.save(str(caminho))
    return caminho


@pytest.fixture
def csv_utf8(tmp_path: Path) -> Path:
    """Arquivo CSV UTF-8 com dados de teste."""
    conteudo = "Nome,RG,Cargo\nAlice,111,Analista\nBob,222,Gerente\n"
    caminho = tmp_path / "dados.csv"
    caminho.write_text(conteudo, encoding="utf-8")
    return caminho


@pytest.fixture
def xlsx_simples(tmp_path: Path) -> Path:
    """Arquivo XLSX com dados de teste."""
    df = pd.DataFrame({"Nome": ["Alice", "Bob"], "RG": ["111", "222"]})
    caminho = tmp_path / "dados.xlsx"
    df.to_excel(caminho, index=False)
    return caminho
